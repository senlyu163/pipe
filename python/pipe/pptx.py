# coding: utf-8

import pptx
from pptx.slide import Slides, Slide
from pptx.chart.data import CategoryDataPoint
import copy

class PPTX(object):
    def __init__(self, prs=None):
        """
        Param prs: a path to a ``.pptx`` file (a string) or a file-like object. 
        """
        self._prs = pptx.Presentation(prs)

    def _illegal_page_num(self, page_num: int) -> bool:
        slides_num = self._prs.slides.__len__()
        if page_num >= slides_num:
            raise IndexError("There have {} slides totally, but get {} index".format(slides_num - 1, page_num))
        return True

    def replace_scatter_data_at_page(self, at_page=None, new_data: list=None) -> bool:
        # suppose there have only one chart in the slide page
        if not new_data or not at_page:
            raise RuntimeError("Please set the page number and new data, got at_page({}) and new_data({})".format(at_page, new_data))

        if self._illegal_page_num(at_page):
            scatter_data = pptx.chart.data.XyChartData()
            series = scatter_data.add_series('')
            series.add_data_point(new_data[0], new_data[1])
            
            shapes = self._prs.slides[at_page].shapes
            for shape in shapes:
                if shape.has_chart:
                    shape.chart.replace_data(scatter_data)
                    return True
                else:
                    continue
        return False

    def slide(self, page_num=0) -> Slide:
        if self._illegal_page_num(page_num):
            return self._prs.slides[page_num]

    def del_slide_page(self, index):
        slides = list(self._prs.slides._sldIdLst)
        self._prs.slides._sldIdLst.remove(slides[index])

    def insert_slide(self, index, new_slide: pptx.oxml.CT_SlideIdList):
        self._prs.slides._sldIdLst.insert(index, new_slide)

    def dul_slide(self, index) -> pptx.oxml.CT_SlideIdList:
        sldLst = list(self._prs.slides._sldIdLst)
        return copy.deepcopy(sldLst[index])
    
    def save(self, file_path):
        self._prs.save(file_path)
    
    @property
    def pages(self):
        return self._prs.slides.__len__()
