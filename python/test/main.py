import sys

sys.path.append("..")

from pipe.pptx import PPTX, pptx, CategoryChartData
from pipe.xlsx import EXCEL
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE


def rep_shape_ctx_retain_fmt(shape, new_context: str = ""):
    shape.text_frame.paragraphs[0].runs[0].text = new_context
    # shape.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT


def first_page(ppt_handle, excel_handle, person_idx: int):
    nm_idx = person_idx + 3
    nm_pos = "A" + str(nm_idx)
    person_name = excel_handle.elements_slice(nm_pos, nm_pos)[0]
    start_pg = person_idx * 4
    page = ppt_handle.slide(start_pg)
    for i, shape in enumerate(page.shapes):
        if isinstance(shape, pptx.shapes.group.GroupShape):
            title_shape = shape.shapes[1]
            rep_shape_ctx_retain_fmt(title_shape, person_name)


def second_page(ppt_handle, excel_handle, person_idx: int):
    nm_idx = person_idx + 3
    scatter_data_start = "B" + str(nm_idx)
    scatter_data_end = "C" + str(nm_idx)
    start_pg = person_idx * 4 + 1

    scatter_data_lst = excel_handle.elements_slice(scatter_data_start, scatter_data_end)
    ppt_handle.replace_scatter_data_at_page(start_pg, scatter_data_lst)

    # replace context
    page = ppt_handle.slide(start_pg)
    for i, shape in enumerate(page.shapes):
        ## text
        if i == 1:
            text_pos = "D" + str(nm_idx)
            shape.text_frame.paragraphs[0].runs[0].text = "管理风格特点： "
            shape.text_frame.paragraphs[0].runs[1].text = excel_handle.elements_slice(text_pos, text_pos)[0]
        ## description
        if i == 9:
            descrip_pos = "E" + str(nm_idx)
            new_ctx = excel_handle.elements_slice(descrip_pos, descrip_pos)[0]
            rep_shape_ctx_retain_fmt(shape, new_ctx)


def third_page(ppt_handle, excel_handle, person_idx: int):
    nm_idx = person_idx + 3
    data_start = "F" + str(nm_idx)
    data_end = "M" + str(nm_idx)
    start_pg = person_idx * 4 + 2

    data_lst = excel_handle.elements_slice(data_start, data_end)
    tag_lst = ['技术职能型', '管理型', '创造创业型', '安全稳定型', '挑战型', '服务型', '生活型', '自主独立型']
    data_dict = {}
    for i in range(len(tag_lst)):
        data_dict[tag_lst[i]] = data_lst[i]

    sorted_data_dict = sorted(data_dict.items(), key=lambda x: x[1], reverse=False)

    sorted_tag_lst = []
    sorted_data_lst = []
    for item in sorted_data_dict:
        sorted_tag_lst.append(item[0])
        sorted_data_lst.append(item[1])

    chart_data = CategoryChartData()
    chart_data.categories = sorted_tag_lst
    chart_data.add_series('', sorted_data_lst)
    # chart_data.categories = ['技术职能型', '管理型', '创造创业型', '安全稳定型', '挑战型', '服务型', '生活型', '自主独立型']
    # chart_data.add_series('', data_lst)
    ppt_handle.replace_chart_data_at_page(start_pg, chart_data)

    # replace context
    page = ppt_handle.slide(start_pg)
    for i, shape in enumerate(page.shapes):
        ## text
        if i == 2:
            text_pos = "N" + str(nm_idx)
            shape.text_frame.paragraphs[0].runs[0].text = excel_handle.elements_slice(text_pos, text_pos)[0]
        # ## description
        if i == 4:
            descrip_pos = "O" + str(nm_idx)
            new_ctx = excel_handle.elements_slice(descrip_pos, descrip_pos)[0]
            # shape.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            # shape.text_frame.word_wrap = True
            shape.text_frame.paragraphs[0].runs[0].text = new_ctx
            shape.text_frame.paragraphs[0].runs[1].text = "" 
            shape.text_frame.paragraphs[0].runs[2].text = "" 
            shape.text_frame.paragraphs[0].runs[3].text = "" 
            shape.text_frame.paragraphs[0].runs[4].text = "" 

        # if shape.has_text_frame:
        #     print("{} : {}".format(i, shape.text))


def fourth_page(ppt_handle, excel_handle, person_idx: int):
    nm_idx = person_idx + 3
    data_start = "P" + str(nm_idx)
    data_end = "S" + str(nm_idx)
    start_pg = person_idx * 4 + 3

    data_lst = excel_handle.elements_slice(data_start, data_end)

    chart_data = CategoryChartData()
    # chart_data.categories = ['感情承诺', '理想承诺', '规范承诺', '机会承诺']
    chart_data.categories = ['机会承诺', '规范承诺', '理想承诺', '感情承诺']
    chart_data.add_series('', data_lst)
    ppt_handle.replace_chart_data_at_page(start_pg, chart_data)

    # replace context
    page = ppt_handle.slide(start_pg)
    for i, shape in enumerate(page.shapes):
        ## text
        if i == 2:
            text_pos = "T" + str(nm_idx)
            shape.text_frame.paragraphs[0].runs[0].text = excel_handle.elements_slice(text_pos, text_pos)[0]


def main():
    excel = EXCEL("raw_data_new.xlsx", False)
    excel.active_ws_by_name(excel.get_sheets_name[1])
    person_num = excel.rows - 2  # sub 2 means dropping the title
    raw = PPTX("raw_ppt.pptx")

    for person_idx in range(person_num):
        # first page
        first_page(raw, excel, person_idx)
        second_page(raw, excel, person_idx)
        third_page(raw, excel, person_idx)
        fourth_page(raw, excel, person_idx)


    raw.save("done.pptx")


if __name__ == '__main__':
    main()